import os
from PyPDF2 import PdfReader
from docx import Document


class FileParser:
    """文件解析服务 - 支持 PDF、DOCX、PPTX、TXT、MD"""

    @staticmethod
    async def parse_pdf(file_path: str) -> str:
        """解析PDF文件"""
        reader = PdfReader(file_path)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n\n".join(texts)

    @staticmethod
    async def parse_docx(file_path: str) -> str:
        """解析DOCX文件"""
        doc = Document(file_path)
        texts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                texts.append(paragraph.text)
        return "\n\n".join(texts)

    @staticmethod
    async def parse_pptx(file_path: str) -> str:
        """解析PPTX文件"""
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(row_text))
            if slide_text:
                text_parts.append(f"=== 第{i}页 ===\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts)

    @staticmethod
    async def parse_txt(file_path: str) -> str:
        """解析TXT文件"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()

    @staticmethod
    async def parse_markdown(file_path: str) -> str:
        """解析Markdown文件"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()

    @staticmethod
    async def parse_file(file_path: str) -> str:
        """自动识别文件类型并解析"""
        ext = os.path.splitext(file_path)[1].lower()
        parsers = {
            ".pdf": FileParser.parse_pdf,
            ".docx": FileParser.parse_docx,
            ".doc": FileParser.parse_docx,
            ".pptx": FileParser.parse_pptx,
            ".ppt": FileParser.parse_pptx,
            ".txt": FileParser.parse_txt,
            ".md": FileParser.parse_markdown,
            ".markdown": FileParser.parse_markdown,
        }
        parser = parsers.get(ext)
        if parser is None:
            # 默认按文本处理
            return await FileParser.parse_txt(file_path)
        return await parser(file_path)


file_parser = FileParser()
